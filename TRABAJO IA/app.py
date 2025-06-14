from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
import io
import openai

# Configura la clave de API de OpenAI
openai.api_key = "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"

# Crea la instancia de Flask
app = Flask(__name__)

# Configura CORS después de inicializar Flask
CORS(app)

# Función para procesar el texto usando GPT de OpenAI
def procesar_con_ia(texto):
    try:
        respuesta = openai.Completion.create(
            engine="text-davinci-003",  # O usa el modelo más adecuado
            prompt=f"Divide este texto en puntos clave para diapositivas:\n\n{texto}",
            max_tokens=500,
            temperature=0.5
        )
        return respuesta.choices[0].text.strip()
    except Exception as e:
        print(f"Error en la llamada a la API de OpenAI: {e}")
        return texto  # En caso de error, devolver el texto sin procesar

@app.route('/generate-slides', methods=['POST'])
def generate_slides():
    data = request.get_json()
    if not data or 'text' not in data:
        return {"error": "Falta el parámetro 'text'"}, 400

    # Procesa el texto usando la IA (OpenAI GPT)
    text = data.get("text", "")
    processed_text = procesar_con_ia(text)  # Llama a la IA para procesar el texto
    slides = processed_text.split("\n\n")  # Divide el texto procesado en bloques

    # Crea la presentación
    prs = Presentation()
    for slide_content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Diapositiva"
        content.text = slide_content

    # Guarda el archivo PPTX en memoria
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    return send_file(pptx_file, as_attachment=True, download_name="diapositivas.pptx")

if __name__ == '__main__':
    app.run(debug=True, port=5001)
